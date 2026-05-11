# -*- coding: utf-8 -*-
"""
Rapor Olusturma Modulu

- IRT ve MCRT analizlerinden tek paket rapor uretir
- Excel ve PPTX dosyalarini olusturur
- Ekranda gorulen ana grafiklerin rapora tasinmasini hedefler
"""

import math
import os
import re
import zipfile
from datetime import datetime
import unicodedata

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
import seaborn as sns
from PIL import Image


PROJE_KOK = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUTPUT_DIR = os.path.join(PROJE_KOK, "output")

SLIDE_W = 13.333
SLIDE_H = 7.5

COLORS = [
    "#6366f1",
    "#10b981",
    "#f59e0b",
    "#ef4444",
    "#8b5cf6",
    "#ec4899",
    "#06b6d4",
    "#14b8a6",
]

THEME = {
    "bg": "0f172a",
    "panel": "111827",
    "card": "ffffff",
    "text": "e5e7eb",
    "muted": "94a3b8",
    "accent": "6366f1",
    "accent2": "8b5cf6",
    "success": "10b981",
    "warning": "f59e0b",
    "danger": "ef4444",
}


def slugify(text):
    text = str(text or "")
    text = (
        text.replace("İ", "I")
        .replace("I", "I")
        .replace("Ş", "S")
        .replace("Ğ", "G")
        .replace("Ü", "U")
        .replace("Ö", "O")
        .replace("Ç", "C")
    )
    text = text.lower()
    text = (
        text.replace("ı", "i")
        .replace("ğ", "g")
        .replace("ü", "u")
        .replace("ş", "s")
        .replace("ö", "o")
        .replace("ç", "c")
    )
    text = unicodedata.normalize("NFKD", text).encode("ascii", "ignore").decode("ascii")
    return re.sub(r"[^a-z0-9]+", "_", text).strip("_")


def _ensure_output_dir():
    os.makedirs(OUTPUT_DIR, exist_ok=True)


def _ifade_kolonu_bul(df):
    for col in df.columns:
        sade = slugify(col)
        if sade in ("ifade", "fade"):
            return col
    raise KeyError("ifade")


def _marka_kolonu_bul(df):
    for col in df.columns:
        sade = slugify(col)
        if sade == "marka":
            return col
    raise KeyError("marka")


def _set_plot_style():
    plt.rcParams["font.family"] = "DejaVu Sans"
    plt.rcParams["axes.unicode_minus"] = False
    sns.set_style("whitegrid")


def _chart_figure(figsize=(12, 7)):
    _set_plot_style()
    fig = plt.figure(figsize=figsize, facecolor="white")
    return fig


def _save_close(fig, path):
    fig.tight_layout()
    fig.savefig(path, dpi=180, bbox_inches="tight", pad_inches=0.12)
    plt.close(fig)


def _nice_axis_upper(max_value, min_upper=20):
    max_value = max(0, float(max_value or 0))
    if max_value <= 0:
        return min_upper
    padded = max_value * 1.18
    if padded <= 20:
        step = 5
    elif padded <= 60:
        step = 10
    else:
        step = 20
    return max(min_upper, int(math.ceil(padded / step) * step))


def _normalize_report_df(df):
    if df is None:
        return pd.DataFrame()
    out = df.copy()
    rename_map = {}
    for col in out.columns:
        sade = slugify(col)
        if sade in ("ifade", "fade"):
            rename_map[col] = "İfade"
        elif sade == "marka":
            rename_map[col] = "Marka"
        elif sade == "explicit":
            rename_map[col] = "Explicit(%)"
        elif sade == "explicit_pct":
            rename_map[col] = "Explicit(%)"
        elif sade == "implicit_guc":
            rename_map[col] = "Implicit_Guc"
        elif sade == "implicit_skor":
            rename_map[col] = "Implicit_Skor"
        elif sade == "mcrt_skor":
            rename_map[col] = "MCRT_Skor"
        elif sade == "secilme_orani":
            rename_map[col] = "Secilme_Orani"
        elif sade == "toplam_secilme":
            rename_map[col] = "Toplam_Secilme"
        elif sade == "ortalama_hiz":
            rename_map[col] = "Ortalama_Hiz"
        elif sade == "medyan_hiz":
            rename_map[col] = "Medyan_Hiz"
        elif sade == "blok":
            rename_map[col] = "Blok"
        elif sade == "kategori":
            rename_map[col] = "Kategori"
        elif sade == "n":
            rename_map[col] = "N"
    out = out.rename(columns=rename_map)
    return out


def _normalize_meta(meta):
    base = {
        "project_name": "",
        "test_type": "IRT",
        "analysis_engine": "IRT",
        "quality": {},
        "correlation": {},
        "statistics": [],
        "category_summary": [],
        "blok_analizleri": {},
        "blok_kaliteleri": {},
    }
    base.update(meta or {})
    return base


def _safe_num(val, default=0):
    try:
        if pd.isna(val):
            return default
        return float(val)
    except Exception:
        return default


def _analysis_mode(meta):
    return "mcrt" if str(meta.get("test_type", "")).lower() in ("mcrt", "mrt") else "irt"


def _metric_columns(df, meta):
    mode = _analysis_mode(meta)
    explicit_col = "Secilme_Orani" if mode == "mcrt" and "Secilme_Orani" in df.columns else "Explicit(%)"
    implicit_candidates = ["MCRT_Skor", "Implicit_Skor", "Implicit_Guc"] if mode == "mcrt" else ["Implicit_Skor", "Implicit_Guc"]
    implicit_col = next((c for c in implicit_candidates if c in df.columns), "Implicit_Guc")
    return explicit_col, implicit_col


def _build_chart_context(df, meta):
    explicit_col, implicit_col = _metric_columns(df, meta)
    marka_col = _marka_kolonu_bul(df)
    ifade_col = _ifade_kolonu_bul(df)
    labels = [str(x) for x in df[ifade_col].dropna().astype(str).unique().tolist()]
    markalar = [str(x) for x in df[marka_col].dropna().astype(str).unique().tolist()]
    return {
        "explicit_col": explicit_col,
        "implicit_col": implicit_col,
        "marka_col": marka_col,
        "ifade_col": ifade_col,
        "labels": labels,
        "markalar": markalar,
    }


def _grouped_values(df, labels, markalar, label_col, marka_col, value_col):
    data = {}
    for marka in markalar:
        row_vals = []
        for label in labels:
            hit = df[(df[label_col].astype(str) == str(label)) & (df[marka_col].astype(str) == str(marka))]
            row_vals.append(_safe_num(hit.iloc[0][value_col], 0) if not hit.empty and value_col in hit.columns else 0)
        data[marka] = row_vals
    return data


def _plot_grouped_bar(df, meta, value_col, title, path, y_min=0, y_max=None):
    ctx = _build_chart_context(df, meta)
    labels = ctx["labels"]
    markalar = ctx["markalar"]
    values = _grouped_values(df, labels, markalar, ctx["ifade_col"], ctx["marka_col"], value_col)
    fig = _chart_figure((13, 7))
    ax = fig.add_subplot(111)
    x = np.arange(len(labels))
    total = max(len(markalar), 1)
    width = 0.8 / total
    for i, marka in enumerate(markalar):
        ax.bar(x - 0.4 + width / 2 + i * width, values[marka], width=width, label=marka, color=COLORS[i % len(COLORS)])
    ax.set_xticks(x)
    label_rotation = 90 if _analysis_mode(meta) == "mcrt" else 20
    label_align = "center" if label_rotation == 90 else "right"
    ax.set_xticklabels(labels, rotation=label_rotation, ha=label_align)
    if y_max is None:
        all_values = [v for vals in values.values() for v in vals]
        y_max = _nice_axis_upper(max(all_values) if all_values else 0, min_upper=20)
    ax.set_ylim(y_min, y_max)
    ax.legend(loc="upper center", ncol=min(4, max(1, len(markalar))), frameon=False)
    ax.grid(axis="y", alpha=0.18)
    _save_close(fig, path)


def _plot_gap(df, meta, path):
    ctx = _build_chart_context(df, meta)
    labels = ctx["labels"]
    markalar = ctx["markalar"]
    explicit_col = ctx["explicit_col"]
    implicit_col = ctx["implicit_col"]
    fig = _chart_figure((13, 7))
    ax = fig.add_subplot(111)
    x = np.arange(len(labels))
    total = max(len(markalar), 1)
    width = 0.8 / total
    gap_rows = {}
    for i, marka in enumerate(markalar):
        values = []
        for label in labels:
            hit = df[(df[ctx["ifade_col"]].astype(str) == str(label)) & (df[ctx["marka_col"]].astype(str) == str(marka))]
            if hit.empty:
                values.append(0)
            else:
                row = hit.iloc[0]
                values.append(_safe_num(row.get(explicit_col), 0) - _safe_num(row.get(implicit_col), 0))
        gap_rows[marka] = values
        ax.bar(x - 0.4 + width / 2 + i * width, values, width=width, label=marka, color=COLORS[i % len(COLORS)])
    ax.axhline(0, color="#475569", linewidth=1)
    ax.set_xticks(x)
    label_rotation = 90 if _analysis_mode(meta) == "mcrt" else 20
    label_align = "center" if label_rotation == 90 else "right"
    ax.set_xticklabels(labels, rotation=label_rotation, ha=label_align)
    max_abs = max([abs(v) for marka_vals in gap_rows.values() for v in marka_vals] or [0])
    gap_upper = _nice_axis_upper(max_abs, min_upper=15)
    ax.set_ylim(-gap_upper, gap_upper)
    ax.legend(loc="upper center", ncol=min(4, max(1, len(markalar))), frameon=False)
    ax.grid(axis="y", alpha=0.18)
    _save_close(fig, path)


def _plot_radar(df, meta, path):
    ctx = _build_chart_context(df, meta)
    labels = ctx["labels"]
    markalar = ctx["markalar"]
    implicit_col = ctx["implicit_col"]
    if not labels:
        labels = ["Veri Yok"]
    angles = np.linspace(0, 2 * math.pi, len(labels), endpoint=False).tolist()
    angles += angles[:1]
    fig = _chart_figure((10, 8))
    ax = fig.add_subplot(111, polar=True)
    for i, marka in enumerate(markalar):
        vals = []
        for label in labels:
            hit = df[(df[ctx["ifade_col"]].astype(str) == str(label)) & (df[ctx["marka_col"]].astype(str) == str(marka))]
            vals.append(_safe_num(hit.iloc[0][implicit_col], 0) if not hit.empty and implicit_col in hit.columns else 0)
        vals += vals[:1]
        color = COLORS[i % len(COLORS)]
        ax.plot(angles, vals, color=color, linewidth=2, label=marka)
        ax.fill(angles, vals, color=color, alpha=0.12)
    ax.set_xticks(angles[:-1])
    ax.set_xticklabels(labels)
    ax.set_ylim(0, 100)
    ax.legend(loc="upper center", bbox_to_anchor=(0.5, -0.08), ncol=min(4, max(1, len(markalar))), frameon=False)
    _save_close(fig, path)


def _plot_scatter(df, meta, path, title=None):
    ctx = _build_chart_context(df, meta)
    explicit_col = ctx["explicit_col"]
    implicit_col = ctx["implicit_col"]
    fig = _chart_figure((13, 7))
    ax = fig.add_subplot(111)
    for i, marka in enumerate(ctx["markalar"]):
        brand_df = df[df[ctx["marka_col"]].astype(str) == str(marka)]
        ax.scatter(
            brand_df[explicit_col].astype(float),
            brand_df[implicit_col].astype(float),
            s=160,
            color=COLORS[i % len(COLORS)],
            label=marka,
            alpha=0.9,
            edgecolors="white",
            linewidths=1.4,
        )
    ax.set_xlim(0, 100)
    ax.set_ylim(0, 100)
    ax.set_xlabel("Secilme Payi %" if _analysis_mode(meta) == "mcrt" else "Explicit Algi %")
    ax.set_ylabel("MCRT Skoru (Dominans)" if _analysis_mode(meta) == "mcrt" else "Implicit Guc / Skor")
    ax.grid(alpha=0.18)
    ax.legend(loc="upper center", ncol=min(4, max(1, len(ctx["markalar"]))), frameon=False)
    _save_close(fig, path)


def _plot_quality_summary(meta, path):
    q = meta.get("quality") or {}
    total = int(q.get("toplam_katilimci") or 0)
    eliminated = int(q.get("elenen_katilimci") or 0)
    valid = int(q.get("kalan_katilimci") or q.get("gecerli_katilimci") or max(total - eliminated, 0))
    avg_speed = _safe_num(q.get("ortalama_hiz") or q.get("ortalama_cevap_suresi"), 0)

    fig = _chart_figure((13, 4.8))
    ax = fig.add_subplot(111)
    ax.axis("off")
    cards = [
        ("TOPLAM KATILIMCI", total, "#1f2937"),
        ("ELENEN KATILIMCI", eliminated, "#ef4444"),
        ("GECERLI KATILIMCI", valid, "#10b981"),
        ("ORT. TEPKI HIZI (MS)", int(avg_speed), "#f59e0b"),
    ]
    for i, (label, val, color) in enumerate(cards):
        x0 = 0.03 + i * 0.24
        rect = plt.Rectangle((x0, 0.18), 0.21, 0.58, facecolor="#f8fafc", edgecolor="#e5e7eb", linewidth=1.2)
        ax.add_patch(rect)
        ax.text(x0 + 0.105, 0.62, label, ha="center", va="center", fontsize=10, color="#64748b", fontweight="bold")
        ax.text(x0 + 0.105, 0.40, str(val), ha="center", va="center", fontsize=24, color=color, fontweight="bold")
    _save_close(fig, path)


def _plot_correlation_summary(meta, path):
    corr = meta.get("correlation") or {}
    pearson = _safe_num(corr.get("pearson"), 0)
    spearman = _safe_num(corr.get("spearman"), 0)

    fig = _chart_figure((13, 4.8))
    ax = fig.add_subplot(111)
    ax.axis("off")
    cards = [
        ("Pearson Korelasyonu", f"{pearson:.3f}", "#6366f1", "Beyan ve bilinçaltı arasındaki doğrusal ilişki gücü."),
        ("Spearman Korelasyonu", f"{spearman:.3f}", "#8b5cf6", "Sıralama tutarlılığı ve marka diziliminin uyumu."),
    ]
    for i, (title, val, color, desc) in enumerate(cards):
        x0 = 0.05 + i * 0.47
        rect = plt.Rectangle((x0, 0.18), 0.40, 0.60, facecolor="#f8fafc", edgecolor=color, linewidth=1.5)
        ax.add_patch(rect)
        ax.text(x0 + 0.20, 0.64, title, ha="center", va="center", fontsize=13, color=color, fontweight="bold")
        ax.text(x0 + 0.20, 0.45, val, ha="center", va="center", fontsize=28, color="#111827", fontweight="bold")
        ax.text(x0 + 0.20, 0.28, desc, ha="center", va="center", fontsize=10, color="#64748b", wrap=True)
    _save_close(fig, path)


def _plot_statistics_table(meta, path):
    rows = meta.get("statistics") or []
    fig = _chart_figure((13, 6.5))
    ax = fig.add_subplot(111)
    ax.axis("off")
    if not rows:
        ax.text(0.5, 0.5, "Bu analiz icin istatistiksel fark tablosu yok.", ha="center", va="center", fontsize=16, color="#64748b")
        _save_close(fig, path)
        return
    cols = ["Karsilastirma", "Explicit Fark", "P-Value (Exp)", "P-Value (Imp)", "Cohen's d", "Yorum"]
    body = []
    for row in rows[:14]:
        body.append([
            f"{row.get('marka_a', '-') } vs {row.get('marka_b', '-')}",
            f"%{row.get('explicit_fark', '-')}",
            f"{row.get('explicit_p', '-')}",
            f"{row.get('implicit_p', '-')}",
            f"{row.get('cohens_d', '-')}",
            f"{row.get('etki_buyuklugu', '-')}",
        ])
    table = ax.table(cellText=body, colLabels=cols, loc="center", cellLoc="center", colLoc="center")
    table.auto_set_font_size(False)
    table.set_fontsize(9)
    table.scale(1, 1.6)
    for (r, c), cell in table.get_celld().items():
        cell.set_edgecolor("#e5e7eb")
        if r == 0:
            cell.set_facecolor("#eef2ff")
            cell.set_text_props(weight="bold", color="#1e1b4b")
        else:
            cell.set_facecolor("#ffffff")
    _save_close(fig, path)


def _plot_metric_table(df, meta, path, title):
    show_cols = [c for c in ["Marka", "İfade", "Explicit(%)", "Implicit_Guc", "Implicit_Skor", "Secilme_Orani", "MCRT_Skor", "Toplam_Secilme", "Blok"] if c in df.columns]
    display_df = df[show_cols].copy().head(16)
    fig = _chart_figure((13, 6.5))
    ax = fig.add_subplot(111)
    ax.axis("off")
    if display_df.empty:
        ax.text(0.5, 0.5, "Gosterilecek veri yok.", ha="center", va="center", fontsize=16, color="#64748b")
        _save_close(fig, path)
        return
    body = display_df.round(2).astype(str).values.tolist()
    table = ax.table(cellText=body, colLabels=list(display_df.columns), loc="center", cellLoc="center", colLoc="center")
    table.auto_set_font_size(False)
    table.set_fontsize(8)
    table.scale(1, 1.45)
    for (r, c), cell in table.get_celld().items():
        cell.set_edgecolor("#e5e7eb")
        if r == 0:
            cell.set_facecolor("#eef2ff")
            cell.set_text_props(weight="bold", color="#1e1b4b")
        else:
            cell.set_facecolor("#ffffff")
    _save_close(fig, path)


def _plot_category_summary(meta, path):
    rows = meta.get("category_summary") or []
    fig = _chart_figure((13, 6.5))
    ax = fig.add_subplot(111)
    ax.axis("off")
    if not rows:
        ax.text(0.5, 0.5, "Kategori ozet verisi yok.", ha="center", va="center", fontsize=16, color="#64748b")
        _save_close(fig, path)
        return
    df = pd.DataFrame(rows).copy().rename(columns={
        "marka": "Marka",
        "kategori": "Kategori",
        "explicit_pct": "Explicit(%)",
        "implicit_skor": "Implicit_Skor",
        "secilme_orani": "Secilme_Orani",
        "mcrt_skor": "MCRT_Skor",
        "n": "N",
    })
    show_cols = [c for c in ["Marka", "Kategori", "Explicit(%)", "Implicit_Skor", "Secilme_Orani", "MCRT_Skor", "N"] if c in df.columns]
    body = df[show_cols].round(2).astype(str).head(18).values.tolist()
    table = ax.table(cellText=body, colLabels=show_cols, loc="center", cellLoc="center", colLoc="center")
    table.auto_set_font_size(False)
    table.set_fontsize(8)
    table.scale(1, 1.45)
    for (r, c), cell in table.get_celld().items():
        cell.set_edgecolor("#e5e7eb")
        if r == 0:
            cell.set_facecolor("#eef2ff")
            cell.set_text_props(weight="bold", color="#1e1b4b")
        else:
            cell.set_facecolor("#ffffff")
    _save_close(fig, path)


def _plot_mcrt_heatmap(df, meta, path):
    ctx = _build_chart_context(df, meta)
    marka_col = ctx["marka_col"]
    ifade_col = ctx["ifade_col"]
    metric_col = "MCRT_Skor" if "MCRT_Skor" in df.columns else ctx["implicit_col"]
    if df.empty or metric_col not in df.columns:
        fig = _chart_figure((13, 7))
        ax = fig.add_subplot(111)
        ax.axis("off")
        ax.text(0.5, 0.5, "Isi haritasi icin veri yok.", ha="center", va="center", fontsize=16, color="#64748b")
        _save_close(fig, path)
        return

    pivot = (
        df[[marka_col, ifade_col, metric_col]]
        .copy()
        .assign(**{metric_col: pd.to_numeric(df[metric_col], errors="coerce").fillna(0)})
        .pivot_table(index=ifade_col, columns=marka_col, values=metric_col, aggfunc="mean", fill_value=0)
    )

    if pivot.empty:
        fig = _chart_figure((13, 7))
        ax = fig.add_subplot(111)
        ax.axis("off")
        ax.text(0.5, 0.5, "Isi haritasi icin veri yok.", ha="center", va="center", fontsize=16, color="#64748b")
        _save_close(fig, path)
        return

    max_rows = 16
    if len(pivot.index) > max_rows:
        row_scores = pivot.max(axis=1).sort_values(ascending=False)
        selected = row_scores.head(max_rows).index.tolist()
        ordered = [idx for idx in pivot.index.tolist() if idx in selected]
        pivot = pivot.loc[ordered]

    fig = _chart_figure((13, 7.4))
    ax = fig.add_subplot(111)
    sns.heatmap(
        pivot,
        cmap=sns.light_palette("#0B6B7F", as_cmap=True),
        linewidths=0.6,
        linecolor="#E5E7EB",
        annot=True,
        fmt=".0f",
        cbar=False,
        ax=ax,
    )
    ax.set_xlabel("")
    ax.set_ylabel("")
    ax.xaxis.tick_top()
    ax.tick_params(axis="x", top=True, bottom=False, labeltop=True, labelbottom=False, labelrotation=0, labelsize=10, pad=6)
    ax.tick_params(axis="y", labelrotation=0, labelsize=10)
    fig.subplots_adjust(top=0.90, bottom=0.04, left=0.19, right=0.98)
    _save_close(fig, path)


def _create_chart_bundle(df, meta, prefix):
    charts = {}
    if df is None or df.empty:
        return charts

    explicit_col, implicit_col = _metric_columns(df, meta)
    charts["quality"] = os.path.join(OUTPUT_DIR, f"{prefix}_quality.png")
    _plot_quality_summary(meta, charts["quality"])

    if _analysis_mode(meta) == "irt":
        charts["correlation"] = os.path.join(OUTPUT_DIR, f"{prefix}_correlation.png")
        _plot_correlation_summary(meta, charts["correlation"])

    charts["scatter"] = os.path.join(OUTPUT_DIR, f"{prefix}_scatter.png")
    _plot_scatter(df, meta, charts["scatter"])

    charts["explicit"] = os.path.join(OUTPUT_DIR, f"{prefix}_explicit.png")
    _plot_grouped_bar(
        df,
        meta,
        explicit_col,
        "Secilme Payi (%)" if _analysis_mode(meta) == "mcrt" else "Explicit Algi (%)",
        charts["explicit"],
    )

    charts["implicit"] = os.path.join(OUTPUT_DIR, f"{prefix}_implicit.png")
    _plot_grouped_bar(
        df,
        meta,
        implicit_col,
        "Zihinsel Dominans (Skor)" if _analysis_mode(meta) == "mcrt" else "Implicit Guc (Skor)",
        charts["implicit"],
    )

    charts["gap"] = os.path.join(OUTPUT_DIR, f"{prefix}_gap.png")
    _plot_gap(df, meta, charts["gap"])

    charts["radar"] = os.path.join(OUTPUT_DIR, f"{prefix}_radar.png")
    _plot_radar(df, meta, charts["radar"])

    charts["table"] = os.path.join(OUTPUT_DIR, f"{prefix}_table.png")
    _plot_metric_table(df, meta, charts["table"], "Analiz Ozet Tablosu")

    charts["category"] = os.path.join(OUTPUT_DIR, f"{prefix}_category.png")
    _plot_category_summary(meta, charts["category"])

    if _analysis_mode(meta) == "irt":
        charts["stats"] = os.path.join(OUTPUT_DIR, f"{prefix}_stats.png")
        _plot_statistics_table(meta, charts["stats"])

    return charts


def _clean_ai_text(ai_text):
    if not ai_text:
        return []
    text = re.sub(r"<br\s*/?>", "\n", ai_text, flags=re.I)
    text = re.sub(r"<h2[^>]*>", "\n## ", text, flags=re.I)
    text = re.sub(r"<h3[^>]*>", "\n### ", text, flags=re.I)
    text = re.sub(r"</(p|div|li|h2|h3|ul)>", "\n", text, flags=re.I)
    text = re.sub(r"<.*?>", "", text)
    text = text.replace("&nbsp;", " ").replace("&quot;", '"')
    lines = [ln.strip() for ln in text.splitlines() if ln.strip()]

    sections = []
    current_title = "Yapay Zeka Uzman Analizi"
    current_body = []

    for line in lines:
        if line.startswith("## "):
            if current_body:
                sections.append((current_title, "\n\n".join(current_body)))
                current_body = []
            current_title = line[3:].strip() or "Yapay Zeka Uzman Analizi"
        elif line.startswith("### "):
            if current_body:
                sections.append((current_title, "\n\n".join(current_body)))
                current_body = []
            current_title = line[4:].strip() or "Yapay Zeka Uzman Analizi"
        else:
            current_body.append(line)

    if current_body:
        sections.append((current_title, "\n\n".join(current_body)))

    chunks = []
    max_chars = 3000
    for title, body in sections:
        if not body.strip():
            continue
        paragraphs = [p.strip() for p in body.split("\n\n") if p.strip()]
        if len(body) <= max_chars:
            chunks.append({"title": title, "body": body})
            continue

        split_index = 1
        active = []
        active_size = 0
        for para in paragraphs:
            para_size = len(para)
            if active and active_size + para_size > max_chars:
                chunks.append({"title": f"{title} ({split_index})", "body": "\n\n".join(active)})
                split_index += 1
                active = [para]
                active_size = para_size
            else:
                active.append(para)
                active_size += para_size
        if active:
            suffix = f" ({split_index})" if split_index > 1 else ""
            chunks.append({"title": f"{title}{suffix}", "body": "\n\n".join(active)})
    return chunks


def _ppt_rgb(value):
    from pptx.dml.color import RGBColor
    return RGBColor.from_string(value)


def _add_slide_title(slide, title, subtitle=None):
    from pptx.util import Inches, Pt
    tx = slide.shapes.add_textbox(Inches(0.7), Inches(0.38), Inches(11.1), Inches(0.92))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(19 if len(str(title or "")) > 58 else 23)
    p.font.bold = True
    p.font.name = "Aptos"
    p.font.color.rgb = _ppt_rgb("0F172A")
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(10)
        p2.font.name = "Aptos"
        p2.font.color.rgb = _ppt_rgb("64748B")


def _prepare_slide(slide, meta=None):
    from pptx.util import Inches, Pt
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = _ppt_rgb("F8FAFC")
    top_band = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(SLIDE_W), Inches(0.18))
    top_band.fill.solid()
    top_band.fill.fore_color.rgb = _ppt_rgb("1E293B")
    top_band.line.fill.background()
    footer = slide.shapes.add_shape(1, Inches(0.7), Inches(7.02), Inches(11.95), Inches(0.01))
    footer.fill.solid()
    footer.fill.fore_color.rgb = _ppt_rgb("E2E8F0")
    footer.line.fill.background()
    if meta:
        tag = slide.shapes.add_textbox(Inches(10.55), Inches(0.42), Inches(2.0), Inches(0.3))
        p = tag.text_frame.paragraphs[0]
        p.text = meta.get("project_name") or (meta.get("test_type") or "").upper()
        p.font.size = Pt(9)
        p.font.name = "Aptos"
        p.alignment = 2
        p.font.color.rgb = _ppt_rgb("64748B")


def _add_note_bar(slide, text):
    from pptx.util import Inches, Pt
    box = slide.shapes.add_shape(1, Inches(0.7), Inches(6.28), Inches(11.95), Inches(0.54))
    box.fill.solid()
    box.fill.fore_color.rgb = _ppt_rgb("EEF2FF")
    box.line.fill.background()
    tx = slide.shapes.add_textbox(Inches(0.9), Inches(6.34), Inches(11.55), Inches(0.42))
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.margin_left = 0
    tf.margin_right = 0
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.name = "Aptos"
    p.font.color.rgb = _ppt_rgb("4338CA")


def _chart_series_matrix(df, meta, value_col):
    ctx = _build_chart_context(df, meta)
    labels = ctx["labels"]
    markalar = ctx["markalar"]
    rows = _grouped_values(df, labels, markalar, ctx["ifade_col"], ctx["marka_col"], value_col)
    return ctx, labels, markalar, rows


def _best_labels(labels, rows, limit=8):
    if len(labels) <= limit:
        return labels, False
    scores = []
    for idx, label in enumerate(labels):
        vals = [abs(_safe_num(values[idx], 0)) for values in rows.values()]
        scores.append((max(vals) if vals else 0, idx, label))
    selected = [label for _, _, label in sorted(scores, reverse=True)[:limit]]
    ordered = [label for label in labels if label in set(selected)]
    return ordered, True


def _add_native_grouped_chart_slide(prs, meta, df, value_col, title, subtitle=None, y_min=0, y_max=100, note=None):
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    from pptx.util import Inches, Pt

    ctx, labels, markalar, rows = _chart_series_matrix(df, meta, value_col)
    labels, truncated = _best_labels(labels, rows, limit=8)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _prepare_slide(slide, meta)
    local_subtitle = subtitle
    if truncated:
        local_subtitle = (subtitle + " | " if subtitle else "") + "En yuksek 8 ifade gosterilmistir."
    _add_slide_title(slide, title, local_subtitle)

    data = CategoryChartData()
    data.categories = labels
    for marka in markalar:
        all_vals = rows.get(marka, [])
        selected_vals = []
        for label in labels:
            idx = ctx["labels"].index(label)
            selected_vals.append(_safe_num(all_vals[idx], 0))
        data.add_series(marka, tuple(selected_vals))

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.85),
        Inches(1.28),
        Inches(11.7),
        Inches(4.85),
        data,
    ).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.TOP
    chart.legend.include_in_layout = False
    chart.chart_style = 10
    chart.value_axis.minimum_scale = y_min
    chart.value_axis.maximum_scale = y_max
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.tick_labels.font.size = Pt(10)
    chart.category_axis.tick_labels.font.size = Pt(10)
    chart.category_axis.tick_labels.font.name = "Aptos"
    chart.value_axis.tick_labels.font.name = "Aptos"
    chart.category_axis.format.line.color.rgb = _ppt_rgb("CBD5E1")
    chart.value_axis.format.line.color.rgb = _ppt_rgb("CBD5E1")
    for i, series in enumerate(chart.series):
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = _ppt_rgb(COLORS[i % len(COLORS)].replace("#", ""))
        series.format.line.color.rgb = _ppt_rgb(COLORS[i % len(COLORS)].replace("#", ""))
    if note:
        _add_note_bar(slide, note)
    return slide


def _add_native_gap_chart_slide(prs, meta, df, title, subtitle=None):
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    from pptx.util import Inches, Pt

    ctx = _build_chart_context(df, meta)
    labels = ctx["labels"]
    markalar = ctx["markalar"]
    explicit_col = ctx["explicit_col"]
    implicit_col = ctx["implicit_col"]
    values_by_marka = {}
    for marka in markalar:
        vals = []
        for label in labels:
            hit = df[(df[ctx["ifade_col"]].astype(str) == str(label)) & (df[ctx["marka_col"]].astype(str) == str(marka))]
            if hit.empty:
                vals.append(0)
            else:
                row = hit.iloc[0]
                vals.append(_safe_num(row.get(explicit_col), 0) - _safe_num(row.get(implicit_col), 0))
        values_by_marka[marka] = vals
    labels, truncated = _best_labels(labels, values_by_marka, limit=8)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _prepare_slide(slide, meta)
    local_subtitle = subtitle
    if truncated:
        local_subtitle = (subtitle + " | " if subtitle else "") + "En belirgin 8 ifade gosterilmistir."
    _add_slide_title(slide, title, local_subtitle)

    data = CategoryChartData()
    data.categories = labels
    for marka in markalar:
        series_vals = []
        for label in labels:
            idx = ctx["labels"].index(label)
            series_vals.append(_safe_num(values_by_marka[marka][idx], 0))
        data.add_series(marka, tuple(series_vals))

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.85),
        Inches(1.28),
        Inches(11.7),
        Inches(4.85),
        data,
    ).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.TOP
    chart.legend.include_in_layout = False
    chart.chart_style = 11
    chart.value_axis.minimum_scale = -100
    chart.value_axis.maximum_scale = 100
    chart.value_axis.crosses_at = 0
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.tick_labels.font.size = Pt(10)
    chart.category_axis.tick_labels.font.size = Pt(10)
    chart.category_axis.format.line.color.rgb = _ppt_rgb("CBD5E1")
    chart.value_axis.format.line.color.rgb = _ppt_rgb("CBD5E1")
    for i, series in enumerate(chart.series):
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = _ppt_rgb(COLORS[i % len(COLORS)].replace("#", ""))
        series.format.line.color.rgb = _ppt_rgb(COLORS[i % len(COLORS)].replace("#", ""))
    _add_note_bar(slide, "Pozitif fark, beyan veya secim payinin hiz skorundan daha baskin oldugunu; negatif fark ise hiz avantajinin daha guclu oldugunu gosterir.")
    return slide


def _add_native_scatter_slide(prs, meta, df, title, subtitle=None):
    from pptx.chart.data import XyChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_MARKER_STYLE
    from pptx.util import Inches, Pt

    ctx = _build_chart_context(df, meta)
    explicit_col = ctx["explicit_col"]
    implicit_col = ctx["implicit_col"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _prepare_slide(slide, meta)
    _add_slide_title(slide, title, subtitle)

    data = XyChartData()
    for marka in ctx["markalar"]:
        series = data.add_series(marka)
        brand_df = df[df[ctx["marka_col"]].astype(str) == str(marka)]
        for _, row in brand_df.iterrows():
            x_val = _safe_num(row.get(explicit_col), 0)
            y_val = _safe_num(row.get(implicit_col), 0)
            series.add_data_point(x_val, y_val)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.XY_SCATTER,
        Inches(0.85),
        Inches(1.28),
        Inches(11.7),
        Inches(4.85),
        data,
    ).chart
    chart.chart_style = 12
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.TOP
    chart.legend.include_in_layout = False
    chart.category_axis.minimum_scale = 0
    chart.category_axis.maximum_scale = 100
    chart.value_axis.minimum_scale = 0
    chart.value_axis.maximum_scale = 100
    chart.category_axis.has_major_gridlines = True
    chart.value_axis.has_major_gridlines = True
    chart.category_axis.tick_labels.font.size = Pt(10)
    chart.value_axis.tick_labels.font.size = Pt(10)
    chart.category_axis.has_title = True
    chart.value_axis.has_title = True
    chart.category_axis.axis_title.text_frame.text = "Secilme Payi %" if _analysis_mode(meta) == "mcrt" else "Explicit Algi %"
    chart.value_axis.axis_title.text_frame.text = "Zihinsel Dominans" if _analysis_mode(meta) == "mcrt" else "Implicit Skor"
    for i, series in enumerate(chart.series):
        series.marker.style = XL_MARKER_STYLE.CIRCLE
        series.marker.size = 9
        series.format.line.color.rgb = _ppt_rgb(COLORS[i % len(COLORS)].replace("#", ""))
        series.marker.format.fill.solid()
        series.marker.format.fill.fore_color.rgb = _ppt_rgb(COLORS[i % len(COLORS)].replace("#", ""))
    _add_note_bar(slide, "Sag ust bolge, hem beyan hem de hiz temelli algida daha guclu sahiplik alanlarini gosterir.")
    return slide


def _add_native_table_slide(prs, meta, title, columns, rows, subtitle=None, note=None, font_size=10):
    from pptx.util import Inches, Pt

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _prepare_slide(slide, meta)
    _add_slide_title(slide, title, subtitle)

    max_rows = min(len(rows), 14)
    rows = rows[:max_rows]
    table_shape = slide.shapes.add_table(len(rows) + 1, len(columns), Inches(0.82), Inches(1.34), Inches(11.75), Inches(4.95))
    table = table_shape.table
    widths = []
    if len(columns) == 6:
        widths = [2.1, 1.2, 1.1, 1.1, 1.0, 3.0]
    elif len(columns) == 7:
        widths = [1.7, 1.6, 1.1, 1.1, 1.1, 1.1, 0.8]
    else:
        even = 11.75 / max(1, len(columns))
        widths = [even] * len(columns)
    for idx, width in enumerate(widths[: len(columns)]):
        table.columns[idx].width = Inches(width)

    for c, col in enumerate(columns):
        cell = table.cell(0, c)
        cell.text = str(col)
        cell.fill.solid()
        cell.fill.fore_color.rgb = _ppt_rgb("E0E7FF")
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(font_size)
        cell.text_frame.paragraphs[0].font.name = "Aptos"
        cell.text_frame.paragraphs[0].font.color.rgb = _ppt_rgb("1E1B4B")

    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = "" if value is None else str(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = _ppt_rgb("FFFFFF" if r % 2 else "F8FAFC")
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(font_size)
            p.font.name = "Aptos"
            p.font.color.rgb = _ppt_rgb("334155")
    if note:
        _add_note_bar(slide, note)
    return slide


def _add_quality_slide(prs, meta):
    from pptx.util import Inches, Pt
    q = meta.get("quality") or {}
    total = int(q.get("toplam_katilimci") or 0)
    eliminated = int(q.get("elenen_katilimci") or 0)
    valid = int(q.get("kalan_katilimci") or q.get("gecerli_katilimci") or max(total - eliminated, 0))
    avg_speed = int(_safe_num(q.get("ortalama_hiz") or q.get("ortalama_cevap_suresi"), 0))

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _prepare_slide(slide, meta)
    _add_slide_title(slide, "Veri Kalitesi ve Orneklem Ozeti", "Sunum akisinda kullanilacak temiz veri tabani")
    cards = [
        ("Toplam Katilimci", total, "1E293B"),
        ("Elenen", eliminated, "DC2626"),
        ("Gecerli", valid, "059669"),
        ("Ort. Tepki Suresi", f"{avg_speed} ms", "D97706"),
    ]
    for i, (label, val, color) in enumerate(cards):
        left = 0.82 + i * 2.95
        box = slide.shapes.add_shape(1, Inches(left), Inches(1.85), Inches(2.55), Inches(2.1))
        box.fill.solid()
        box.fill.fore_color.rgb = _ppt_rgb("FFFFFF")
        box.line.color.rgb = _ppt_rgb("E2E8F0")
        kpi = slide.shapes.add_textbox(Inches(left + 0.18), Inches(2.18), Inches(2.15), Inches(0.55))
        p = kpi.text_frame.paragraphs[0]
        p.text = str(val)
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.name = "Aptos"
        p.font.color.rgb = _ppt_rgb(color)
        lbl = slide.shapes.add_textbox(Inches(left + 0.18), Inches(2.88), Inches(2.15), Inches(0.36))
        p2 = lbl.text_frame.paragraphs[0]
        p2.text = label
        p2.font.size = Pt(11)
        p2.font.name = "Aptos"
        p2.font.color.rgb = _ppt_rgb("64748B")
    _add_note_bar(slide, "Bu sayfa, raporda kullanilan orneklemin temizlik ve hiz profilini tek bakista verir.")
    return slide


def _add_correlation_slide(prs, meta):
    from pptx.util import Inches, Pt
    corr = meta.get("correlation") or {}
    pearson = _safe_num(corr.get("pearson"), 0)
    spearman = _safe_num(corr.get("spearman"), 0)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _prepare_slide(slide, meta)
    _add_slide_title(slide, "Explicit ve Implicit Tutarlilik", "Beyan edilen algi ile hiz temelli algi arasindaki uyum")
    cards = [
        ("Pearson", f"{pearson:.3f}", "4F46E5", "Dogrusallik"),
        ("Spearman", f"{spearman:.3f}", "7C3AED", "Siralama uyumu"),
    ]
    for i, (label, val, color, foot) in enumerate(cards):
        left = 1.15 + i * 5.9
        box = slide.shapes.add_shape(1, Inches(left), Inches(1.9), Inches(4.55), Inches(2.25))
        box.fill.solid()
        box.fill.fore_color.rgb = _ppt_rgb("FFFFFF")
        box.line.color.rgb = _ppt_rgb(color)
        box.line.width = Pt(1.4)
        tx = slide.shapes.add_textbox(Inches(left + 0.25), Inches(2.22), Inches(4.0), Inches(0.42))
        p = tx.text_frame.paragraphs[0]
        p.text = label
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.name = "Aptos"
        p.font.color.rgb = _ppt_rgb(color)
        tx2 = slide.shapes.add_textbox(Inches(left + 0.25), Inches(2.72), Inches(4.0), Inches(0.72))
        p2 = tx2.text_frame.paragraphs[0]
        p2.text = val
        p2.font.size = Pt(28)
        p2.font.bold = True
        p2.font.name = "Aptos"
        p2.font.color.rgb = _ppt_rgb("0F172A")
        tx3 = slide.shapes.add_textbox(Inches(left + 0.25), Inches(3.38), Inches(4.0), Inches(0.28))
        p3 = tx3.text_frame.paragraphs[0]
        p3.text = foot
        p3.font.size = Pt(10)
        p3.font.name = "Aptos"
        p3.font.color.rgb = _ppt_rgb("64748B")
    _add_note_bar(slide, "Korelasyonun pozitif ve yuksek olmasi, acik beyan ile hiz temelli algi arasinda daha tutarli bir desen oldugunu gosterir.")
    return slide


def _add_picture_slide(prs, title, image_path, subtitle=None, note=None):
    from pptx.util import Inches
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _prepare_slide(slide)
    _add_slide_title(slide, title, subtitle=subtitle)

    box_left = 0.75
    box_top = 1.45
    box_width = 11.85
    box_height = 5.2

    with Image.open(image_path) as img:
        img_w, img_h = img.size

    if not img_w or not img_h:
        slide.shapes.add_picture(image_path, Inches(box_left), Inches(box_top), width=Inches(box_width), height=Inches(box_height))
        return slide

    img_ratio = img_w / img_h
    box_ratio = box_width / box_height

    if img_ratio >= box_ratio:
        final_width = box_width
        final_height = box_width / img_ratio
        final_left = box_left
        final_top = box_top + (box_height - final_height) / 2
    else:
        final_height = box_height
        final_width = box_height * img_ratio
        final_top = box_top
        final_left = box_left + (box_width - final_width) / 2

    slide.shapes.add_picture(
        image_path,
        Inches(final_left),
        Inches(final_top),
        width=Inches(final_width),
        height=Inches(final_height),
    )
    if note:
        _add_note_bar(slide, note)
    return slide


def _add_cover_slide(prs, meta):
    from pptx.util import Inches, Pt
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = _ppt_rgb("F8FAFC")
    accent = slide.shapes.add_shape(1, Inches(0.7), Inches(0.9), Inches(4.0), Inches(5.8))
    accent.fill.solid()
    accent.fill.fore_color.rgb = _ppt_rgb("0F172A")
    accent.line.fill.background()

    tx = slide.shapes.add_textbox(Inches(1.05), Inches(1.35), Inches(3.2), Inches(3.8))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = meta.get("project_name") or "Marka Algi Raporu"
    p.font.size = Pt(27)
    p.font.bold = True
    p.font.name = "Aptos"
    p.font.color.rgb = _ppt_rgb("FFFFFF")

    p2 = tf.add_paragraph()
    p2.text = f"Analiz Motoru: {meta.get('analysis_engine') or meta.get('test_type')}"
    p2.font.size = Pt(16)
    p2.font.name = "Aptos"
    p2.font.color.rgb = _ppt_rgb("CBD5E1")

    p3 = tf.add_paragraph()
    p3.text = datetime.now().strftime("%d.%m.%Y %H:%M")
    p3.font.size = Pt(12)
    p3.font.name = "Aptos"
    p3.font.color.rgb = _ppt_rgb("94A3B8")

    strap = slide.shapes.add_textbox(Inches(4.95), Inches(1.52), Inches(6.9), Inches(1.6))
    p4 = strap.text_frame.paragraphs[0]
    p4.text = "Karar odakli, okunur ve yonetime sunulabilir marka algi raporu"
    p4.font.size = Pt(24)
    p4.font.bold = True
    p4.font.name = "Aptos Display"
    p4.font.color.rgb = _ppt_rgb("0F172A")
    p5 = strap.text_frame.add_paragraph()
    p5.text = "Bu sunum, ekran goruntusu mantigindan uzaklasmis; grafik, tablo ve karar cümlesi merkezli bir danismanlik raporu dilinde kurgulanmistir."
    p5.font.size = Pt(12)
    p5.font.name = "Aptos"
    p5.font.color.rgb = _ppt_rgb("475569")
    meta_box = slide.shapes.add_shape(1, Inches(4.95), Inches(4.75), Inches(5.2), Inches(1.1))
    meta_box.fill.solid()
    meta_box.fill.fore_color.rgb = _ppt_rgb("FFFFFF")
    meta_box.line.color.rgb = _ppt_rgb("E2E8F0")
    meta_tx = slide.shapes.add_textbox(Inches(5.2), Inches(5.02), Inches(4.7), Inches(0.65))
    p6 = meta_tx.text_frame.paragraphs[0]
    p6.text = f"Test Turu: {(meta.get('test_type') or '').upper()}   |   Surum: Kurumsal Rapor"
    p6.font.size = Pt(12)
    p6.font.name = "Aptos"
    p6.font.color.rgb = _ppt_rgb("334155")


def _add_ai_slides(prs, ai_text, max_slides=3):
    chunks = _clean_ai_text(ai_text)
    if not chunks:
        return
    chunks = chunks[:max_slides]  # Slayt sayisini sinirla
    from pptx.util import Inches, Pt
    for idx, chunk in enumerate(chunks, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _prepare_slide(slide)
        title = chunk.get("title") or (f"Yapay Zeka Uzman Analizi {idx}" if len(chunks) > 1 else "Yapay Zeka Uzman Analizi")
        _add_slide_title(slide, title, "Yapay zeka uzman degerlendirmesi")
        box = slide.shapes.add_textbox(Inches(0.9), Inches(1.55), Inches(11.15), Inches(5.1))
        tf = box.text_frame
        tf.word_wrap = True
        paragraphs = [p.strip() for p in str(chunk.get("body") or "").split("\n\n") if p.strip()]
        for p_idx, para in enumerate(paragraphs):
            p = tf.paragraphs[0] if p_idx == 0 else tf.add_paragraph()
            p.text = para
            p.font.size = Pt(12)
            p.font.name = "Aptos"
            p.font.color.rgb = _ppt_rgb("1F2937")
            p.space_after = Pt(8)


def _statistics_rows(meta):
    rows = []
    for row in (meta.get("statistics") or [])[:12]:
        rows.append([
            f"{row.get('marka_a', '-') } vs {row.get('marka_b', '-')}",
            row.get("explicit_fark", "-"),
            row.get("explicit_p", "-"),
            row.get("implicit_p", "-"),
            row.get("cohens_d", "-"),
            row.get("etki_buyuklugu", "-"),
        ])
    return rows


def _metric_table_rows(df, meta):
    if df is None or df.empty:
        return [], []

    ctx = _build_chart_context(df, meta)
    ifade_col = ctx["ifade_col"]
    marka_col = ctx["marka_col"]
    explicit_col = ctx["explicit_col"]
    implicit_col = ctx["implicit_col"]

    rows = []
    for ifade in ctx["labels"][:12]:
        grp = df[df[ifade_col].astype(str) == str(ifade)].copy()
        if grp.empty:
            continue
        explicit_winner = grp.sort_values(explicit_col, ascending=False).iloc[0]
        implicit_winner = grp.sort_values(implicit_col, ascending=False).iloc[0]
        explicit_gap = _safe_num(explicit_winner.get(explicit_col), 0) - _safe_num(grp[explicit_col].astype(float).mean(), 0)
        implicit_gap = _safe_num(implicit_winner.get(implicit_col), 0) - _safe_num(grp[implicit_col].astype(float).mean(), 0)

        yorum = []
        if str(explicit_winner[marka_col]) == str(implicit_winner[marka_col]):
            yorum.append("iki katmanda da lider")
        else:
            yorum.append("beyan ve hiz lideri ayrisiyor")
        if explicit_gap >= 10 or implicit_gap >= 10:
            yorum.append("fark belirgin")
        else:
            yorum.append("fark sinirli")

        rows.append([
            str(ifade),
            str(explicit_winner[marka_col]),
            f"{_safe_num(explicit_winner.get(explicit_col), 0):.0f}",
            str(implicit_winner[marka_col]),
            f"{_safe_num(implicit_winner.get(implicit_col), 0):.0f}",
            ", ".join(yorum),
        ])

    return ["Ifade", "Explicit Lider", "Skor", "Implicit Lider", "Skor", "Yorum"], rows


def _category_table_rows(meta):
    rows = meta.get("category_summary") or []
    if not rows:
        return [], []
    df = pd.DataFrame(rows).copy().rename(columns={
        "marka": "Marka",
        "kategori": "Kategori",
        "explicit_pct": "Explicit(%)",
        "implicit_skor": "Implicit_Skor",
        "secilme_orani": "Secilme_Orani",
        "mcrt_skor": "MCRT_Skor",
        "n": "N",
    })
    show_cols = [c for c in ["Marka", "Kategori", "Explicit(%)", "Implicit_Skor", "Secilme_Orani", "MCRT_Skor", "N"] if c in df.columns]
    return show_cols, df[show_cols].round(2).astype(str).head(12).values.tolist()


def _pptx_olustur_final(yol, meta, ana_charts, blok_chartlari, ai_text):
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    _add_cover_slide(prs, meta)
    ordered = [
        ("Veri Kalitesi ve Temizlik Ozeti", "quality"),
        ("Bilimsel Korelasyon Skorlari", "correlation"),
        ("Bilincalti Algi Haritasi (Scatter Plot)" if _analysis_mode(meta) == "irt" else "Zihinsel Harita (MCRT)", "scatter"),
        ("Explicit Algi (%)" if _analysis_mode(meta) == "irt" else "Secilme Payi (%)", "explicit"),
        ("Implicit Guc (Skor)" if _analysis_mode(meta) == "irt" else "Zihinsel Dominans (Skor)", "implicit"),
        ("GAP Analizi (Beyan vs Bilincalti)" if _analysis_mode(meta) == "irt" else "Pay vs Hiz Analizi", "gap"),
        ("Marka Algi Agi (Radar)" if _analysis_mode(meta) == "irt" else "Stimulus Bazli Dagilim (Radar)", "radar"),
        ("Istatistiksel Farklar", "stats"),
        ("Kategori Ozeti", "category"),
        ("Analiz Ozet Tablosu", "table"),
    ]

    for title, key in ordered:
        if key in ana_charts and os.path.exists(ana_charts[key]):
            _add_picture_slide(prs, title, ana_charts[key], subtitle=meta.get("project_name"))

    for blok_adi, charts in (blok_chartlari or {}).items():
        for title, key in ordered:
            if key in charts and os.path.exists(charts[key]):
                _add_picture_slide(
                    prs,
                    f"{title} - {blok_adi.replace('_', ' ').title()}",
                    charts[key],
                    subtitle=meta.get("project_name"),
                )

    _add_ai_slides(prs, ai_text)
    prs.save(yol)


def _plot_mcrt_heatmap_v2(df, meta, path):
    ctx = _build_chart_context(df, meta)
    marka_col = ctx["marka_col"]
    ifade_col = ctx["ifade_col"]
    metric_col = "MCRT_Skor" if "MCRT_Skor" in df.columns else ctx["implicit_col"]

    fig = _chart_figure((13, 7.4))
    ax = fig.add_subplot(111)
    ax.axis("off")

    if df.empty or metric_col not in df.columns:
        ax.text(0.5, 0.5, "Isi haritasi icin veri yok.", ha="center", va="center", fontsize=16, color="#64748b")
        _save_close(fig, path)
        return

    work_df = df[[marka_col, ifade_col, metric_col]].copy()
    work_df[metric_col] = pd.to_numeric(work_df[metric_col], errors="coerce").fillna(0)
    pivot = work_df.pivot_table(index=ifade_col, columns=marka_col, values=metric_col, aggfunc="mean", fill_value=0)

    if pivot.empty:
        ax.text(0.5, 0.5, "Isi haritasi icin veri yok.", ha="center", va="center", fontsize=16, color="#64748b")
        _save_close(fig, path)
        return

    ax.clear()
    sns.heatmap(
        pivot,
        cmap=sns.light_palette("#0B6B7F", as_cmap=True),
        linewidths=0.6,
        linecolor="#E5E7EB",
        annot=True,
        fmt=".0f",
        cbar=False,
        ax=ax,
    )
    ax.set_xlabel("")
    ax.set_ylabel("")
    ax.tick_params(axis="x", labelrotation=0, labelsize=10)
    ax.tick_params(axis="y", labelrotation=0, labelsize=10)
    _save_close(fig, path)


def _create_mcrt_heatmap_pages(df, meta, prefix, rows_per_page=24):
    if df is None or df.empty:
        return []

    ctx = _build_chart_context(df, meta)
    marka_col = ctx["marka_col"]
    ifade_col = ctx["ifade_col"]
    metric_col = "MCRT_Skor" if "MCRT_Skor" in df.columns else ctx["implicit_col"]
    if metric_col not in df.columns:
        return []

    work_df = df[[marka_col, ifade_col, metric_col]].copy()
    work_df[metric_col] = pd.to_numeric(work_df[metric_col], errors="coerce").fillna(0)
    pivot = work_df.pivot_table(index=ifade_col, columns=marka_col, values=metric_col, aggfunc="mean", fill_value=0)
    if pivot.empty:
        return []

    page_paths = []
    total_pages = max(1, math.ceil(len(pivot.index) / rows_per_page))
    for page_index in range(total_pages):
        start = page_index * rows_per_page
        end = start + rows_per_page
        chunk = pivot.iloc[start:end].copy()
        path = os.path.join(OUTPUT_DIR, f"{prefix}_heatmap_page_{page_index + 1}.png")
        fig = _chart_figure((13, 7.4))
        ax = fig.add_subplot(111)
        row_count = len(chunk.index)
        y_size = 9 if row_count <= 16 else 8
        x_size = 10 if row_count <= 16 else 9
        annot_size = 11 if row_count <= 14 else (10 if row_count <= 18 else 9)
        sns.heatmap(
            chunk,
            cmap=sns.light_palette("#0B6B7F", as_cmap=True),
            linewidths=0.6,
            linecolor="#E5E7EB",
            annot=True,
            fmt=".0f",
            cbar=False,
            annot_kws={"size": annot_size},
            ax=ax,
        )
        ax.set_xlabel("")
        ax.set_ylabel("")
        ax.xaxis.tick_top()
        ax.tick_params(axis="x", top=True, bottom=False, labeltop=True, labelbottom=False, labelrotation=0, labelsize=x_size, pad=6)
        ax.tick_params(axis="y", labelrotation=0, labelsize=y_size)
        fig.subplots_adjust(top=0.90, bottom=0.04, left=0.19, right=0.98)
        _save_close(fig, path)
        page_paths.append(path)
    return page_paths


def _plot_metric_table_v2(df, meta, path, title):
    if _analysis_mode(meta) == "mcrt":
        preferred_cols = ["Blok", "Marka", "Ä°fade", "Secilme_Orani", "MCRT_Skor", "Ortalama_Hiz", "Toplam_Secilme"]
        row_limit = 18
    else:
        preferred_cols = ["Marka", "Ä°fade", "Explicit(%)", "Implicit_Guc", "Implicit_Skor"]
        row_limit = 16

    show_cols = [c for c in preferred_cols if c in df.columns]
    display_df = df[show_cols].copy().head(row_limit)
    fig = _chart_figure((13, 6.5))
    ax = fig.add_subplot(111)
    ax.axis("off")
    if display_df.empty:
        ax.text(0.5, 0.5, "Gosterilecek veri yok.", ha="center", va="center", fontsize=16, color="#64748b")
        _save_close(fig, path)
        return
    body = display_df.round(2).astype(str).values.tolist()
    table = ax.table(cellText=body, colLabels=list(display_df.columns), loc="center", cellLoc="center", colLoc="center")
    table.auto_set_font_size(False)
    table.set_fontsize(8)
    table.scale(1, 1.45)
    for (r, c), cell in table.get_celld().items():
        cell.set_edgecolor("#e5e7eb")
        if r == 0:
            cell.set_facecolor("#eef2ff")
            cell.set_text_props(weight="bold", color="#1e1b4b")
        else:
            cell.set_facecolor("#ffffff")
    _save_close(fig, path)


def _mcrt_slide_note(key, df, meta):
    if df is None or df.empty:
        return "Bu slayt, mevcut raporda yorumlanabilir veri bulunmadigini gosterir."

    ctx = _build_chart_context(df, meta)
    marka_col = ctx["marka_col"]
    ifade_col = ctx["ifade_col"]
    explicit_col = ctx["explicit_col"]
    implicit_col = ctx["implicit_col"]

    work = df.copy()
    for col in [explicit_col, implicit_col]:
        if col in work.columns:
            work[col] = pd.to_numeric(work[col], errors="coerce").fillna(0)

    if key == "heatmap":
        top_row = work.sort_values(implicit_col, ascending=False).iloc[0]
        top_brands = work.groupby(marka_col)[implicit_col].mean().sort_values(ascending=False).head(2).index.tolist()
        top_expr = (
            work.groupby(ifade_col)[implicit_col]
            .mean()
            .sort_values(ascending=False)
            .head(2)
            .index.tolist()
        )
        return (
            f"Bu tablo, ifade bazinda hangi markalarin daha sistematik sahiplik kurdugunu gosterir. "
            f"Ozellikle {', '.join(map(str, top_expr))} gibi ifadelerde {', '.join(map(str, top_brands))} daha guclu bir yogunluk uretiyor; "
            f"bu da rekabetin rastgele degil, belirli anlam alanlarinda toplandigini gosteriyor."
        )

    if key == "explicit":
        marka_ort = work.groupby(marka_col)[explicit_col].mean().sort_values(ascending=False)
        leader = marka_ort.index[0]
        top_expr = (
            work.groupby(ifade_col)[explicit_col]
            .mean()
            .sort_values(ascending=False)
            .head(3)
            .index.tolist()
        )
        return (
            f"Bu grafik, secilme payi uzerinden markalarin hangi ifadelerde daha kolay sahiplik kurdugunu gosterir. "
            f"En guclu acik sahiplik deseni {leader} tarafinda gorulurken, ozellikle {', '.join(map(str, top_expr))} gibi ifadeler toplam secimi daha fazla tasiyor."
        )

    if key == "implicit":
        marka_ort = work.groupby(marka_col)[implicit_col].mean().sort_values(ascending=False)
        leader = marka_ort.index[0]
        leading_expr = (
            work[work[marka_col].astype(str) == str(leader)]
            .sort_values(implicit_col, ascending=False)[ifade_col]
            .astype(str)
            .head(3)
            .tolist()
        )
        return (
            f"Bu grafik, hiz ve secilme birlikteliginden dogan zihinsel dominansi gosterir. "
            f"{leader} markasi ozellikle {', '.join(map(str, leading_expr))} gibi ifadelerde daha akici ve dogal bir bag kuruyor; "
            f"bu da tercihin sadece secilme degil, zihinsel kolaylik boyutunda da desteklendigini gosteriyor."
        )

    if key == "gap":
        gaps = work[explicit_col] - work[implicit_col]
        largest = work.iloc[gaps.abs().idxmax()]
        return (
            f"Bu grafik, secilme payi ile hiz arasindaki farki gosterir. "
            f"En belirgin ayrisma {largest[marka_col]} icin '{largest[ifade_col]}' ifadesinde gorulmustur; "
            f"bu alan, markanin acik beyanla sahiplenildigi ama otomatik karar akisinin ayni gucte desteklemedigi bir gerilim noktasi olabilir."
        )

    if key == "stats":
        stats_rows = meta.get("statistics") or []
        if not stats_rows:
            return "Bu tablo, markalar arasinda istatistiksel olarak anlamli bir ayrisma bulunmadigini ya da yeterli karsilastirma verisi olmadigini gosterir."
        first = stats_rows[0]
        return (
            f"Bu tablo, markalar arasindaki ayrismanin istatistiksel olarak ne kadar guclu oldugunu gosterir. "
            f"En dikkat cekici karsilastirma {first.get('marka_a', '-')} ve {first.get('marka_b', '-')} arasinda goruluyor; "
            f"p-degerleri ve etki buyuklugu, farkin tesadufi olup olmadigini yorumlamak icin kullanilir."
        )

    if key == "category":
        cat_rows = pd.DataFrame(meta.get("category_summary") or [])
        if cat_rows.empty:
            return "Bu tablo, kategori bazinda anlamli bir sonuc uretmek icin yeterli veri bulunmadigini gosterir."
        metric = "implicit_skor" if "implicit_skor" in cat_rows.columns else ("mcrt_skor" if "mcrt_skor" in cat_rows.columns else None)
        if metric:
            cat_rows[metric] = pd.to_numeric(cat_rows[metric], errors="coerce").fillna(0)
            top = cat_rows.sort_values(metric, ascending=False).iloc[0]
            low = cat_rows.sort_values(metric, ascending=True).iloc[0]
            return (
                f"Bu tablo, kategoriler bazinda markalarin toplu performansini gosterir. "
                f"En guclu kategori eslesmesi {top.get('marka', '-')} icin '{top.get('kategori', '-')}', en zayif eslesme ise {low.get('marka', '-')} icin '{low.get('kategori', '-')}' olarak okunur."
            )
        return "Bu tablo, kategori bazinda markalarin toplu performansini gosterir ve hangi anlam alanlarinin onde oldugunu okumaya yarar."

    if key == "quality":
        q = meta.get("quality") or {}
        total = int(_safe_num(q.get("toplam_katilimci"), 0))
        valid = int(_safe_num(q.get("gecerli_katilimci"), total))
        excluded = int(_safe_num(q.get("elenen_katilimci"), max(total - valid, 0)))
        avg_rt = int(_safe_num(q.get("ortalama_rt_ms"), q.get("ortalama_tepki_ms", 0)))
        return (
            f"Bu slayt, veri setinin ne kadar temiz ve kullanilabilir oldugunu gosterir. "
            f"Toplam {total} katilimcinin {valid} tanesi analize dahil edilmis, {excluded} kayit kalite nedeniyle disarida birakilmistir; ortalama tepki hizi {avg_rt} ms seviyesindedir."
        )

    return "Bu slayt, ilgili metrikteki yuksek ve dusuk alanlari gostererek markanin nerede guclu, nerede desteklenmeye ihtiyac duydugunu okumaya yardimci olur."


def _plot_expression_focus_v2(df, meta, path, ifade_label):
    ctx = _build_chart_context(df, meta)
    marka_col = ctx["marka_col"]
    ifade_col = ctx["ifade_col"]
    explicit_col = ctx["explicit_col"]
    implicit_col = ctx["implicit_col"]
    hiz_col = "Ortalama_Hiz" if "Ortalama_Hiz" in df.columns else ("Medyan_Hiz" if "Medyan_Hiz" in df.columns else None)

    work = df[df[ifade_col].astype(str) == str(ifade_label)].copy()
    fig = _chart_figure((13, 6.7))
    ax = fig.add_subplot(111)

    if work.empty:
        ax.axis("off")
        ax.text(0.5, 0.5, "Bu ifade icin veri yok.", ha="center", va="center", fontsize=16, color="#64748b")
        _save_close(fig, path)
        return

    work[explicit_col] = pd.to_numeric(work[explicit_col], errors="coerce").fillna(0)
    work[implicit_col] = pd.to_numeric(work[implicit_col], errors="coerce").fillna(0)
    if hiz_col:
        work[hiz_col] = pd.to_numeric(work[hiz_col], errors="coerce").fillna(0)
    if marka_col in work.columns:
        agg_map = {
            explicit_col: "mean",
            implicit_col: "mean",
        }
        if hiz_col:
            agg_map[hiz_col] = "mean"
        work = work.groupby(marka_col, as_index=False).agg(agg_map)
    work = work.sort_values(implicit_col, ascending=True)

    y = np.arange(len(work))
    ax.barh(y - 0.18, work[explicit_col], height=0.32, color="#0B6B7F", label="Secilme Payi %")
    ax.barh(y + 0.18, work[implicit_col], height=0.32, color="#1D4ED8", label="MCRT Skor")
    ax.set_yticks(y)
    ax.set_yticklabels(work[marka_col].astype(str).tolist())
    ax.set_xlim(0, max(20, _nice_axis_upper(max(work[explicit_col].max(), work[implicit_col].max()), min_upper=20)))
    ax.grid(axis="x", alpha=0.18)
    ax.legend(loc="upper right", frameon=False)

    for idx, (_, row) in enumerate(work.iterrows()):
        if hiz_col and row.get(hiz_col, 0) > 0:
            ax.text(
                min(ax.get_xlim()[1] - 1, max(row[explicit_col], row[implicit_col]) + 1.0),
                idx,
                f"{int(row[hiz_col])} ms",
                va="center",
                fontsize=9,
                color="#475569",
            )

    _save_close(fig, path)


def _mcrt_expression_note(df, meta, ifade_label):
    ctx = _build_chart_context(df, meta)
    marka_col = ctx["marka_col"]
    ifade_col = ctx["ifade_col"]
    explicit_col = ctx["explicit_col"]
    implicit_col = ctx["implicit_col"]
    hiz_col = "Ortalama_Hiz" if "Ortalama_Hiz" in df.columns else ("Medyan_Hiz" if "Medyan_Hiz" in df.columns else None)

    work = df[df[ifade_col].astype(str) == str(ifade_label)].copy()
    if work.empty:
        return f"Bu slayt, '{ifade_label}' ifadesi icin anlamli veri bulunmadigini gosterir."

    for col in [explicit_col, implicit_col]:
        work[col] = pd.to_numeric(work[col], errors="coerce").fillna(0)
    if hiz_col:
        work[hiz_col] = pd.to_numeric(work[hiz_col], errors="coerce").fillna(0)
    if marka_col in work.columns:
        agg_map = {
            explicit_col: "mean",
            implicit_col: "mean",
        }
        if hiz_col:
            agg_map[hiz_col] = "mean"
        work = work.groupby(marka_col, as_index=False).agg(agg_map)

    top_score = work.sort_values(implicit_col, ascending=False).iloc[0]
    low_score = work.sort_values(implicit_col, ascending=True).iloc[0]
    top_share = work.sort_values(explicit_col, ascending=False).iloc[0]

    note = (
        f"Bu slayt, '{ifade_label}' ifadesinde tum markalarin birlikte nasil ayrildigini gosterir. "
        f"En yuksek zihinsel dominans {top_score[marka_col]} markasinda, en yuksek secilme payi ise {top_share[marka_col]} markasinda gorulmustur."
    )
    if str(top_score[marka_col]) != str(top_share[marka_col]):
        note += " Bu ayrisma, beyan edilen tercih ile otomatik akisin farkli bir marka lehine calistigini gosterir."
    if hiz_col and top_score.get(hiz_col, 0) > 0:
        note += f" En dusuk performans ise {low_score[marka_col]} tarafinda kalmis; hiz verisi {int(top_score[hiz_col])} ms bandinda daha akici bir karar desenine isaret eder."
    return note


def _ai_slide_note(meta, key, blok_adi=None, ifade_label=None):
    pack = (meta.get("ai_slide_pack") or {}) if isinstance(meta, dict) else {}
    if not pack:
        return None
    if ifade_label and blok_adi:
        expr = ((pack.get("expression_notes") or {}).get(blok_adi) or {}).get(ifade_label)
        if expr:
            return str(expr)
    note_map = pack.get("slide_notes") or {}
    if blok_adi:
        scoped = note_map.get(f"{blok_adi}.{key}")
        if scoped:
            return str(scoped)
    plain = note_map.get(key)
    return str(plain) if plain else None


def _mcrt_expression_focus_charts(df, meta, prefix):
    charts = {}
    if df is None or df.empty:
        return charts
    ctx = _build_chart_context(df, meta)
    for ifade in ctx["labels"]:
        key = slugify(str(ifade)) or "ifade"
        path = os.path.join(OUTPUT_DIR, f"{prefix}_{key}_focus.png")
        _plot_expression_focus_v2(df, meta, path, ifade)
        charts[str(ifade)] = path
    return charts


def _create_full_expression_tables(df, meta, prefix, rows_per_page=18):
    if df is None or df.empty:
        return []

    work = df.copy()
    ctx = _build_chart_context(work, meta)
    ifade_col = ctx["ifade_col"]
    marka_col = ctx["marka_col"]

    sort_cols = [col for col in [ifade_col, marka_col] if col in work.columns]
    if sort_cols:
        work = work.sort_values(sort_cols).reset_index(drop=True)

    preferred_cols = []
    if "Blok" in work.columns:
        preferred_cols.append("Blok")
    preferred_cols.extend([
        ifade_col,
        marka_col,
        ctx["explicit_col"],
        ctx["implicit_col"],
    ])
    for extra in ["Ortalama_Hiz", "Toplam_Secilme", "N"]:
        if extra in work.columns:
            preferred_cols.append(extra)

    display_cols = []
    seen = set()
    for col in preferred_cols:
        if col in work.columns and col not in seen:
            display_cols.append(col)
            seen.add(col)

    page_paths = []
    total_rows = len(work)
    total_pages = max(1, math.ceil(total_rows / rows_per_page))
    for page_index in range(total_pages):
        start = page_index * rows_per_page
        end = start + rows_per_page
        chunk = work.iloc[start:end][display_cols].copy()
        path = os.path.join(OUTPUT_DIR, f"{prefix}_all_expr_{page_index + 1}.png")

        fig = _chart_figure((13, 6.5))
        ax = fig.add_subplot(111)
        ax.axis("off")

        body = chunk.round(2).astype(str).values.tolist()
        table = ax.table(
            cellText=body,
            colLabels=list(chunk.columns),
            loc="center",
            cellLoc="center",
            colLoc="center",
        )
        table.auto_set_font_size(False)
        table.set_fontsize(8)
        table.scale(1, 1.42)
        for (r, c), cell in table.get_celld().items():
            cell.set_edgecolor("#e5e7eb")
            if r == 0:
                cell.set_facecolor("#eef2ff")
                cell.set_text_props(weight="bold", color="#1e1b4b")
            else:
                cell.set_facecolor("#ffffff")
        _save_close(fig, path)
        page_paths.append(path)

    return page_paths


def _irt_expression_scatter_charts(df, meta, prefix):
    charts = {}
    if df is None or df.empty:
        return charts
    ctx = _build_chart_context(df, meta)
    for ifade in ctx["labels"]:
        filtered = df[df[ctx["ifade_col"]].astype(str) == str(ifade)].copy()
        if filtered.empty:
            continue
        key = slugify(str(ifade)) or "ifade"
        path = os.path.join(OUTPUT_DIR, f"{prefix}_{key}_scatter.png")
        _plot_scatter(filtered, meta, path)
        charts[str(ifade)] = path
    return charts


def _irt_expression_note(df, meta, ifade_label):
    ctx = _build_chart_context(df, meta)
    marka_col = ctx["marka_col"]
    ifade_col = ctx["ifade_col"]
    explicit_col = ctx["explicit_col"]
    implicit_col = ctx["implicit_col"]
    work = df[df[ifade_col].astype(str) == str(ifade_label)].copy()
    if work.empty:
        return f"Bu slayt, '{ifade_label}' ifadesi icin anlamli veri olmadigini gosterir."
    work[explicit_col] = pd.to_numeric(work[explicit_col], errors="coerce").fillna(0)
    work[implicit_col] = pd.to_numeric(work[implicit_col], errors="coerce").fillna(0)
    top_exp = work.sort_values(explicit_col, ascending=False).iloc[0]
    top_imp = work.sort_values(implicit_col, ascending=False).iloc[0]
    note = (
        f"Bu slayt, '{ifade_label}' ifadesinde markalarin beyan ve bilinçalti performansini birlikte gosterir. "
        f"Explicit tarafta en yuksek deger {top_exp[marka_col]}, implicit tarafta ise {top_imp[marka_col]} markasindadir."
    )
    if str(top_exp[marka_col]) != str(top_imp[marka_col]):
        note += " Bu ayrisma, tüketicinin söyledigi ile otomatik erisim deseninin farkli markalari öne cikardigini gösterir."
    return note


def _fallback_mcrt_ai_text(df, meta):
    if df is None or df.empty:
        return ""
    ctx = _build_chart_context(df, meta)
    explicit_col = ctx["explicit_col"]
    implicit_col = ctx["implicit_col"]
    marka_col = ctx["marka_col"]
    ifade_col = ctx["ifade_col"]
    work = df.copy()
    work[explicit_col] = pd.to_numeric(work[explicit_col], errors="coerce").fillna(0)
    work[implicit_col] = pd.to_numeric(work[implicit_col], errors="coerce").fillna(0)

    marka_mean = work.groupby(marka_col)[implicit_col].mean().sort_values(ascending=False)
    if marka_mean.empty:
        return ""
    top_brand = str(marka_mean.index[0])
    low_brand = str(marka_mean.index[-1])
    top_row = work.sort_values(implicit_col, ascending=False).iloc[0]
    tension_row = work.assign(_gap=(work[explicit_col] - work[implicit_col]).abs()).sort_values("_gap", ascending=False).iloc[0]

    lines = [
        "Yonetici Ozeti",
        f"Genel tablo, zihinsel dominans tarafinda en guclu markanin {top_brand}, en zayif markanin ise {low_brand} oldugunu gosteriyor.",
        "",
        "Ana Bulgular",
        f"En dikkat cekici ifade-marka eslesmesi {top_row[marka_col]} icin '{top_row[ifade_col]}' alaninda goruluyor.",
        f"Beyan ile hiz arasindaki en belirgin ayrisma ise {tension_row[marka_col]} markasinda '{tension_row[ifade_col]}' ifadesinde ortaya cikiyor.",
        "",
        "Onemli Oneriler",
        f"{top_brand} icin mevcut guclu ifade alanlari iletişim ve yaratıcı içerikte daha görünür hale getirilmeli.",
        f"{low_brand} icin zayif kalan alanlarda ürün deneyimi, raf anlatısı veya mesaj mimarisi yeniden ele alinmali.",
        "",
        "Bu arastirma bize gosteriyor ki",
        "Marka performansini anlamak icin secilme payi ile karar akisina birlikte bakmak gerekiyor; en saglam aksiyon alanlari bu iki sinyalin birlikte guclu oldugu ifadelerde ortaya cikiyor.",
    ]
    return "\n".join(lines)


def _create_chart_bundle_v2(df, meta, prefix):
    charts = {}
    if df is None or df.empty:
        return charts

    explicit_col, implicit_col = _metric_columns(df, meta)
    charts["quality"] = os.path.join(OUTPUT_DIR, f"{prefix}_quality.png")
    _plot_quality_summary(meta, charts["quality"])

    if _analysis_mode(meta) == "irt":
        charts["correlation"] = os.path.join(OUTPUT_DIR, f"{prefix}_correlation.png")
        _plot_correlation_summary(meta, charts["correlation"])

    charts["scatter"] = os.path.join(OUTPUT_DIR, f"{prefix}_scatter.png")
    _plot_scatter(df, meta, charts["scatter"])

    charts["explicit"] = os.path.join(OUTPUT_DIR, f"{prefix}_explicit.png")
    _plot_grouped_bar(
        df,
        meta,
        explicit_col,
        "Secilme Payi (%)" if _analysis_mode(meta) == "mcrt" else "Explicit Algi (%)",
        charts["explicit"],
    )

    charts["implicit"] = os.path.join(OUTPUT_DIR, f"{prefix}_implicit.png")
    _plot_grouped_bar(
        df,
        meta,
        implicit_col,
        "Zihinsel Dominans (Skor)" if _analysis_mode(meta) == "mcrt" else "Implicit Guc (Skor)",
        charts["implicit"],
    )

    charts["gap"] = os.path.join(OUTPUT_DIR, f"{prefix}_gap.png")
    _plot_gap(df, meta, charts["gap"])

    if _analysis_mode(meta) == "mcrt":
        charts["heatmap"] = os.path.join(OUTPUT_DIR, f"{prefix}_heatmap.png")
        _plot_mcrt_heatmap_v2(df, meta, charts["heatmap"])

    charts["radar"] = os.path.join(OUTPUT_DIR, f"{prefix}_radar.png")
    _plot_radar(df, meta, charts["radar"])

    charts["table"] = os.path.join(OUTPUT_DIR, f"{prefix}_table.png")
    _plot_metric_table_v2(df, meta, charts["table"], "Analiz Ozet Tablosu")

    charts["category"] = os.path.join(OUTPUT_DIR, f"{prefix}_category.png")
    _plot_category_summary(meta, charts["category"])

    if _analysis_mode(meta) == "irt" or (meta.get("statistics") or []):
        charts["stats"] = os.path.join(OUTPUT_DIR, f"{prefix}_stats.png")
        _plot_statistics_table(meta, charts["stats"])

    return charts



def _add_executive_summary_slide(prs, meta):
    """Yonetici Ozeti slaytı - AI slide pack'ten cekilen 3-4 madde."""
    from pptx.util import Inches, Pt

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _prepare_slide(slide, meta)
    _add_slide_title(slide, "Yonetici Ozeti", "Bu arastirmadan cikan en kritik noktalar")

    pack = meta.get("ai_slide_pack") or {}
    items = pack.get("yonetici_ozeti") or []
    if not items:
        items = ["Yonetici ozeti icin AI analizi bekleniyor."]

    start_y = 1.85
    for i, item in enumerate(items[:4]):
        # Bullet number circle
        num_box = slide.shapes.add_shape(
            9,  # oval
            Inches(1.0),
            Inches(start_y + i * 1.12),
            Inches(0.45),
            Inches(0.45),
        )
        num_box.fill.solid()
        num_box.fill.fore_color.rgb = _ppt_rgb("4F46E5")
        num_box.line.fill.background()
        num_tx = num_box.text_frame
        num_tx.margin_top = 0
        num_tx.margin_bottom = 0
        p = num_tx.paragraphs[0]
        p.text = str(i + 1)
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.name = "Aptos"
        p.font.color.rgb = _ppt_rgb("FFFFFF")
        p.alignment = 1  # center

        # Bullet text
        tx = slide.shapes.add_textbox(
            Inches(1.65),
            Inches(start_y + i * 1.12),
            Inches(9.8),
            Inches(0.9),
        )
        tf = tx.text_frame
        tf.word_wrap = True
        p2 = tf.paragraphs[0]
        p2.text = str(item)
        p2.font.size = Pt(14)
        p2.font.name = "Aptos"
        p2.font.color.rgb = _ppt_rgb("1E293B")

    _add_note_bar(slide, "Bu ozet, tum raporun en kritik cikarimlarini tek bakista verir.")
    return slide


def _pptx_olustur_final_v2(yol, meta, ana_charts, blok_chartlari, ai_text, blok_frames=None, blok_ifade_charts=None, irt_ifade_scatter_charts=None, appendix_tables=None, heatmap_pages=None, blok_heatmap_pages=None):
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    # 1. Kapak
    _add_cover_slide(prs, meta)

    # 2. Yonetici Ozeti (AI slide pack'ten)
    _add_executive_summary_slide(prs, meta)

    # 3. Ana grafikler
    if _analysis_mode(meta) == "mcrt":
        ordered = [
            ("MCRT Istatistiksel Karsilastirma (Marka Ayrismasi)", "stats"),
            ("Isi Haritasi (MCRT)", "heatmap"),
            ("Secilme Payi (%)", "explicit"),
            ("Zihinsel Dominans (Skor)", "implicit"),
            ("Pay vs Hiz Analizi", "gap"),
        ]
    else:
        ordered = [
            ("Istatistiksel Farklar", "stats"),
            ("Explicit Algi (%)", "explicit"),
            ("Implicit Guc (Skor)", "implicit"),
            ("GAP Analizi (Beyan vs Bilincalti)", "gap"),
            ("Marka Algi Agi (Radar)", "radar"),
        ]

    for title, key in ordered:
        if _analysis_mode(meta) == "mcrt" and key == "heatmap" and (heatmap_pages or []):
            total_pages = len(heatmap_pages or [])
            for page_index, image_path in enumerate(heatmap_pages or [], start=1):
                note = _ai_slide_note(meta, key) or _mcrt_slide_note(key, meta.get("_sonuc_df"), meta)
                page_title = title if total_pages == 1 else f"{title} ({page_index}/{total_pages})"
                _add_picture_slide(prs, page_title, image_path, subtitle=meta.get("project_name"), note=note)
            continue
        if key in ana_charts and os.path.exists(ana_charts[key]):
            if _analysis_mode(meta) == "mcrt":
                note = _ai_slide_note(meta, key) or _mcrt_slide_note(key, meta.get("_sonuc_df"), meta)
            else:
                note = None
            _add_picture_slide(prs, title, ana_charts[key], subtitle=meta.get("project_name"), note=note)

    # 4. Blok grafikleri - SADECE ISI HARITASI (tekrar eden bar grafikleri kaldirildi)
    for blok_adi, charts in (blok_chartlari or {}).items():
        blok_df = (blok_frames or {}).get(blok_adi)
        blok_meta = dict(meta)
        blok_meta["project_name"] = meta.get("project_name")

        if _analysis_mode(meta) == "mcrt" and ((blok_heatmap_pages or {}).get(blok_adi) or []):
            pages = (blok_heatmap_pages or {}).get(blok_adi) or []
            total_pages = len(pages)
            for page_index, image_path in enumerate(pages, start=1):
                page_title = f"Isi Haritasi - {blok_adi.replace('_', ' ').title()}" if total_pages == 1 else f"Isi Haritasi - {blok_adi.replace('_', ' ').title()} ({page_index}/{total_pages})"
                _add_picture_slide(
                    prs,
                    page_title,
                    image_path,
                    subtitle=meta.get("project_name"),
                    note=_ai_slide_note(blok_meta, "heatmap", blok_adi=blok_adi) or _mcrt_slide_note("heatmap", blok_df, blok_meta),
                )
        elif "heatmap" in charts and os.path.exists(charts["heatmap"]):
            _add_picture_slide(
                prs,
                f"Isi Haritasi - {blok_adi.replace('_', ' ').title()}",
                charts["heatmap"],
                subtitle=meta.get("project_name"),
                note=_ai_slide_note(blok_meta, "heatmap", blok_adi=blok_adi) or _mcrt_slide_note("heatmap", blok_df, blok_meta) if _analysis_mode(meta) == "mcrt" else None,
            )

        # Ifade-basina-slayt KALDIRILDI - gereksiz tekrar

    # 5. IRT appendix tablolari (varsa)
    for appendix in (appendix_tables or []):
        image_path = appendix.get("path")
        if image_path and os.path.exists(image_path):
            _add_picture_slide(
                prs,
                appendix.get("title") or "Tum Ifadeler",
                image_path,
                subtitle=meta.get("project_name"),
                note=None,
            )

    # 6. AI Uzman Analizi - EN FAZLA 3 SLAYT
    _add_ai_slides(prs, ai_text, max_slides=3)
    prs.save(yol)


def rapor_olustur(sonuc_df, marka_df, profil_df=None, ai_text=None, rapor_meta=None):
    _ensure_output_dir()
    meta = _normalize_meta(rapor_meta)
    sonuc_df = _normalize_report_df(sonuc_df)
    marka_df = _normalize_report_df(marka_df)
    meta["_sonuc_df"] = sonuc_df

    tarih_str = datetime.now().strftime("%Y%m%d_%H%M%S")
    dosyalar = {}

    excel_dosya = os.path.join(OUTPUT_DIR, f"rapor_{tarih_str}.xlsx")
    with pd.ExcelWriter(excel_dosya, engine="openpyxl") as writer:
        sonuc_df.to_excel(writer, sheet_name="Analiz_Verileri", index=False)
        marka_df.to_excel(writer, sheet_name="Marka_Skorlari", index=False)
        kategori_df = pd.DataFrame(meta.get("category_summary") or [])
        if not kategori_df.empty:
            kategori_df.to_excel(writer, sheet_name="Kategori_Ozeti", index=False)
        if profil_df is not None and not getattr(profil_df, "empty", True):
            profil_df.to_excel(writer, sheet_name="Katilimci_Profilleri", index=False)
    dosyalar["excel"] = excel_dosya

    ana_charts = _create_chart_bundle_v2(sonuc_df, meta, f"ana_{tarih_str}")
    for key, path in ana_charts.items():
        dosyalar[f"chart_{key}"] = path
    if _analysis_mode(meta) == "mcrt":
        heatmap_pages = _create_mcrt_heatmap_pages(sonuc_df, meta, f"ana_{tarih_str}", rows_per_page=24)
        for idx, path in enumerate(heatmap_pages, start=1):
            dosyalar[f"ana_heatmap_page_{idx}"] = path

    blok_chartlari = {}
    blok_frames = {}
    blok_ifade_charts = {}
    irt_ifade_scatter_charts = {}
    appendix_tables = []
    heatmap_pages = []
    blok_heatmap_pages = {}

    if _analysis_mode(meta) == "irt":
        irt_ifade_scatter_charts = _irt_expression_scatter_charts(sonuc_df, meta, f"ana_{tarih_str}")
        for ifade, path in irt_ifade_scatter_charts.items():
            dosyalar[f"{slugify(ifade)}_scatter"] = path
        ana_appendix = _create_full_expression_tables(sonuc_df, meta, f"ana_{tarih_str}", rows_per_page=18)
        for idx, path in enumerate(ana_appendix, start=1):
            dosyalar[f"ana_appendix_{idx}"] = path
            appendix_tables.append({
                "title": f"Tum Ifadeler - Tam Liste ({idx}/{len(ana_appendix)})",
                "path": path,
                "note": "Bu ek tablo, IRT tarafinda tum ifadeleri eksiksiz gormek icin eklenmistir.",
            })

    for blok_adi, kayitlar in (meta.get("blok_analizleri") or {}).items():
        blok_df = _normalize_report_df(pd.DataFrame(kayitlar or []))
        if blok_df.empty:
            continue
        blok_frames[blok_adi] = blok_df
        blok_meta = dict(meta)
        blok_meta["blok_analizleri"] = {}
        blok_meta["quality"] = (meta.get("blok_kaliteleri") or {}).get(blok_adi) or meta.get("quality") or {}
        charts = _create_chart_bundle_v2(blok_df, blok_meta, f"{slugify(blok_adi)}_{tarih_str}")
        blok_chartlari[blok_adi] = charts
        for key, path in charts.items():
            dosyalar[f"{slugify(blok_adi)}_{key}"] = path

        if _analysis_mode(meta) == "mcrt":
            heatmap_page_list = _create_mcrt_heatmap_pages(blok_df, blok_meta, f"{slugify(blok_adi)}_{tarih_str}", rows_per_page=24)
            blok_heatmap_pages[blok_adi] = heatmap_page_list
            for idx, path in enumerate(heatmap_page_list, start=1):
                dosyalar[f"{slugify(blok_adi)}_heatmap_page_{idx}"] = path

            ifade_charts = _mcrt_expression_focus_charts(blok_df, blok_meta, f"{slugify(blok_adi)}_{tarih_str}")
            blok_ifade_charts[blok_adi] = ifade_charts
            for ifade, path in ifade_charts.items():
                dosyalar[f"{slugify(blok_adi)}_{slugify(ifade)}_focus"] = path

            # MCRT deck'te tum ifadeler appendix'i raporu gereksiz uzatiyor.
            # Tum ifadeler analizi ekranda ve ifade odak slaytlarinda zaten gorulebildigi icin
            # bu ham tablo appendix'ini MCRT PPT akisina eklemiyoruz.

    if _analysis_mode(meta) == "mcrt" and not ai_text:
        ai_text = _fallback_mcrt_ai_text(sonuc_df, meta)

    pptx_dosya = os.path.join(OUTPUT_DIR, f"rapor_{tarih_str}.pptx")
    _pptx_olustur_final_v2(
        pptx_dosya,
        meta,
        ana_charts,
        blok_chartlari,
        ai_text,
        blok_frames=blok_frames,
        blok_ifade_charts=blok_ifade_charts,
        irt_ifade_scatter_charts=irt_ifade_scatter_charts,
        appendix_tables=appendix_tables,
        heatmap_pages=heatmap_pages,
        blok_heatmap_pages=blok_heatmap_pages,
    )
    dosyalar["pptx"] = pptx_dosya

    return dosyalar


def rapor_paketi_olustur(dosyalar, paket_on_eki="rapor_paketi"):
    _ensure_output_dir()
    tarih_str = datetime.now().strftime("%Y%m%d_%H%M%S")
    paket_yolu = os.path.join(OUTPUT_DIR, f"{paket_on_eki}_{tarih_str}.zip")
    with zipfile.ZipFile(paket_yolu, "w", compression=zipfile.ZIP_DEFLATED) as zf:
        for _, tam_yol in (dosyalar or {}).items():
            if tam_yol and os.path.exists(tam_yol):
                zf.write(tam_yol, arcname=os.path.basename(tam_yol))
    return paket_yolu
