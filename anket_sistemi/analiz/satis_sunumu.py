# -*- coding: utf-8 -*-
"""
Marka Odaklı Stratejik Sunum Hazırlayıcı (V2 - Senaryo Bazlı)
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Dosya yolları
PROJE_KOK = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
ARTIFACTS_DIR = r"C:\Users\neuropc\.gemini\antigravity\brain\ebe18c0e-9bb6-4b9e-89b7-dbe426e1aeed"
OUTPUT_PATH = os.path.join(PROJE_KOK, "output", "Marka_Strateji_IRT_Sunumu.pptx")

# Görsel yolları
COVER_IMG = os.path.join(ARTIFACTS_DIR, "irt_sales_cover_1777285830886.png")
METHOD_IMG = os.path.join(ARTIFACTS_DIR, "irt_methodology_visual_1777285843682.png")
DASHBOARD_IMG = os.path.join(ARTIFACTS_DIR, "irt_dashboard_visual_premium_1777285858306.png")

def sunum_hazirla():
    os.makedirs(os.path.dirname(OUTPUT_PATH), exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # --- SLAYT 1: KAPAK (VİZYON) ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if os.path.exists(COVER_IMG):
        slide.shapes.add_picture(COVER_IMG, 0, 0, width=prs.slide_width)
    
    left, top, width, height = Inches(0.5), Inches(2.5), Inches(12), Inches(2)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "MARKANIZIN GERÇEK GÜCÜNÜ ÖLÇÜN"
    p.font.bold = True; p.font.size = Pt(48); p.font.color.rgb = RGBColor(255, 255, 255)
    
    p2 = tf.add_paragraph()
    p2.text = "Bilinçaltı Analizi ile Stratejik Karar Destek Mekanizması"
    p2.font.size = Pt(26); p2.font.color.rgb = RGBColor(220, 220, 220)

    # --- SLAYT 2: NEDEN ŞİMDİ? ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Müşteriniz Size Ne Söylemiyor?"
    tf = slide.placeholders[1].text_frame
    tf.text = "Tüketiciler rasyonel olduklarını düşünürler, ancak kararları duygusaldır.\n\n- 'Kaliteli' diyorlar ama gerçekten 'Güveniyorlar' mı?\n- 'Pahalı' diyorlar ama içten içe 'Vazgeçilmez' mi buluyorlar?\n- Rakiplerinizi mi tercih ediyorlar yoksa sadece 'Alışkanlık' mı?"

    # --- SLAYT 3: SENARYO 1 - LANSMAN BAŞARISI ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Vaka 1: Yeni Ürün/Ambalaj Lansmanı"
    tf = slide.placeholders[1].text_frame
    tf.text = "Yeni tasarımınız raflarda 'Heyecan' mı yaratacak yoksa 'Karışıklık' mı?\n\n- Geleneksel anketler 'Güzel' der geçer.\n- IRT Analizi, yeni tasarımın tüketicide 'Premium' algısını milisaniyeler içinde tetikleyip tetiklemediğini ölçer.\n- Riskleri lansman öncesi minimize edin."

    # --- SLAYT 4: SENARYO 2 - FİYAT ESNEKLİĞİ ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Vaka 2: Fiyat Artışı ve Marka Sadakati"
    tf = slide.placeholders[1].text_frame
    tf.text = "Markanız bir fiyat artışına ne kadar hazır?\n\n- Tüketicinin bilinçaltında markanız 'Lüks' kategorisinde mi yoksa 'Fiyat Odaklı' mı?\n- Eğer bilinçaltı bağı güçlüyse, rasyonel fiyat tepkilerini absorbe edebilirsiniz.\n- IRT ile 'Gizli Sadakati' ölçün, fiyat stratejinizi veriye dayandırın."

    # --- SLAYT 5: SENARYO 3 - RAKİP SAVAŞLARI ---
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Vaka 3: Rakipten Pay Çalma (Benchmarking)"
    if os.path.exists(METHOD_IMG):
        slide.shapes.add_picture(METHOD_IMG, Inches(0.5), Inches(1.5), width=Inches(5))
    
    left, top, width, height = Inches(6), Inches(2), Inches(6.5), Inches(4)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "Rakibinizin Zayıf Karnını Bulun"
    p.font.bold = True; p.font.size = Pt(22); p.font.color.rgb = RGBColor(99, 102, 241)
    tf.add_paragraph().text = "Rakip marka 'Samimi' algısında beyan bazında güçlü olsa da, bilinçaltında 'Zayıf Bağ' kurmuş olabilir. Bu boşluğu görerek iletişim stratejinizi rakibin en zayıf olduğu bilinçaltı noktasına kurun."

    # --- SLAYT 6: STRATEJİK HARİTA (SCATTER PLOT ANLAMI) ---
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Stratejik Algı Haritası Size Ne Anlatır?"
    if os.path.exists(DASHBOARD_IMG):
        slide.shapes.add_picture(DASHBOARD_IMG, Inches(4.5), Inches(1.8), width=Inches(8))
    
    left, top, width, height = Inches(0.5), Inches(2), Inches(4), Inches(4)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "● GİZLİ FIRSAT: İnsanların henüz dile dökemediği ama içten içe sevdiği alanlar.\n● STRATEJİK GÜÇ: Markanızın kalesi olan, sarsılmaz algı noktaları.\n● RASYONEL ALGI: Sadece mantıklı olduğu için tercih edildiğiniz riskli bölgeler."

    # --- SLAYT 7: NEDEN BİZ? (İŞ DEĞERİ) ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Karar Vericiler İçin Net Çıktılar"
    tf = slide.placeholders[1].text_frame
    tf.text = "✓ Hız: 1 hafta yerine 48 saatte veri toplama.\n✓ Kesinlik: Sosyal beğeni hatasından arındırılmış %100 dürüst veri.\n✓ Aksiyon: 'Veri' değil, 'İletişim Stratejisi' teslim ediyoruz.\n✓ AI Desteği: Sonuçlar üzerinde saatlerce düşünmeyin, AI uzman yorumlarını kullanın."

    # --- SLAYT 8: İLETİŞİM ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Geleceği Milisaniyelerle Yönetin"
    tf = slide.placeholders[1].text_frame
    tf.text = "\n\nMarkanızın bilinçaltı karnesini bugün birlikte çıkaralım."
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    prs.save(OUTPUT_PATH)
    print(f"Stratejik sunum olusturuldu: {OUTPUT_PATH}")

if __name__ == "__main__":
    sunum_hazirla()
